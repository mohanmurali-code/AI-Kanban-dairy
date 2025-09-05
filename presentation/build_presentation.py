import os
import textwrap
import requests
import urllib.parse
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

OUTPUT_PPTX = "/workspace/presentation/Craniofacial_Syndromes_Anesthesia.pptx"
IMAGES_DIR = "/workspace/presentation/images"

os.makedirs(IMAGES_DIR, exist_ok=True)

# Syndromes list with description and anesthesia points
syndromes = [
    {
        "name": "Apert Syndrome",
        "description": "Craniosynostosis with turribrachycephaly, midface hypoplasia, syndactyly; possible cleft palate.",
        "anesthesia": [
            "Anticipate difficult mask/intubation (midface hypoplasia, choanal atresia)",
            "Screen for CHD; manage ICP if present",
            "Eye protection for proptosis; careful positioning"
        ],
        "image_url": None,
        "image_caption": "Apert syndrome: syndactyly (Wikimedia Commons)"
    },
    {
        "name": "Crouzon Syndrome",
        "description": "Craniosynostosis, midface hypoplasia, beaked nose, proptosis; no limb anomalies.",
        "anesthesia": [
            "Potential difficult airway; consider fiberoptic",
            "Monitor for increased ICP",
            "Eye protection due to proptosis"
        ],
        "image_url": None,
        "image_caption": "Crouzon syndrome (Wikimedia Commons)"
    },
    {
        "name": "Pfeiffer Syndrome",
        "description": "Craniosynostosis with midface hypoplasia; broad, medially deviated thumbs/toes.",
        "anesthesia": [
            "Airway obstruction risk; prepare for difficult intubation",
            "Assess for CHD and increased ICP",
            "Careful positioning and eye protection"
        ],
        "image_url": None,
        "image_caption": "Pfeiffer syndrome (Wikimedia Commons)"
    },
    {
        "name": "Saethre-Chotzen Syndrome",
        "description": "Craniosynostosis, facial asymmetry, ptosis, mild syndactyly.",
        "anesthesia": [
            "Airway usually manageable; assess individually",
            "Monitor for ICP and hearing issues",
            "IV access may be challenging with limb anomalies"
        ],
        "image_url": None,
        "image_caption": "Saethre-Chotzen syndrome (Wikimedia Commons)"
    },
    {
        "name": "Carpenter Syndrome",
        "description": "Craniosynostosis, syndactyly/polydactyly, obesity; developmental delay.",
        "anesthesia": [
            "Difficult ventilation/intubation (midface hypoplasia, obesity)",
            "Cardiac defects common; evaluate preop",
            "Positioning challenges due to habitus and limbs"
        ],
        "image_url": None,
        "image_caption": "Carpenter syndrome (Wikimedia Commons)"
    },
    {
        "name": "Muenke Syndrome",
        "description": "Coronal craniosynostosis, brachycephaly, midface hypoplasia, hearing loss.",
        "anesthesia": [
            "Prepare for potential difficult airway",
            "Assess hearing/communication needs",
            "Monitor for signs of ICP"
        ],
        "image_url": None,
        "image_caption": "Muenke syndrome (Wikimedia Commons)"
    },
    {
        "name": "Treacher Collins Syndrome",
        "description": "Mandibulofacial dysostosis: mandibular/malar hypoplasia, downward palpebral fissures, ear anomalies; cleft palate common.",
        "anesthesia": [
            "Very difficult airway; fiberoptic or video techniques",
            "Consider postoperative airway obstruction; extended monitoring",
            "Address hearing impairment communication"
        ],
        "image_url": None,
        "image_caption": "Treacher Collins syndrome (Wikimedia Commons)"
    },
    {
        "name": "Pierre Robin Sequence",
        "description": "Triad: micrognathia, glossoptosis, cleft palate; significant airway obstruction.",
        "anesthesia": [
            "Maintain spontaneous ventilation until airway secured",
            "Prepare for difficult intubation; alternative airways",
            "Postop monitoring for obstruction"
        ],
        "image_url": None,
        "image_caption": "Pierre Robin sequence (Wikimedia Commons)"
    },
    {
        "name": "Hemifacial Microsomia / Goldenhar",
        "description": "Unilateral facial underdevelopment; ear anomalies; vertebral +/- cardiac defects.",
        "anesthesia": [
            "Difficult intubation due to asymmetry/mandibular hypoplasia",
            "Evaluate cardiac/vertebral anomalies",
            "Careful positioning and padding"
        ],
        "image_url": None,
        "image_caption": "Hemifacial microsomia (Wikimedia Commons)"
    },
    {
        "name": "Freeman-Sheldon Syndrome",
        "description": "Whistling face: microstomia, pursed lips; limb contractures; scoliosis.",
        "anesthesia": [
            "Extremely difficult airway; microstomia",
            "Positioning challenges; consider MH risk per institutional protocol",
            "Plan postoperative respiratory support"
        ],
        "image_url": None,
        "image_caption": "Freeman-Sheldon syndrome (Wikimedia Commons)"
    },
]


def download_image(url: str, dest_path: str) -> str | None:
    try:
        headers = {
            "User-Agent": (
                "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 "
                "(KHTML, like Gecko) Chrome/123.0 Safari/537.36"
            ),
            "Referer": "https://commons.wikimedia.org/",
            "Accept": "image/avif,image/webp,image/apng,image/*,*/*;q=0.8",
        }
        resp = requests.get(url, timeout=30, headers=headers)
        resp.raise_for_status()
        with open(dest_path, "wb") as f:
            f.write(resp.content)
        return dest_path
    except Exception as e:
        print(f"Failed to download {url}: {e}")
        return None


def add_title_slide(prs: Presentation, title: str, subtitle: str):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_syndrome_slide(prs: Presentation, item: dict, image_path: str | None):
    slide_layout = prs.slide_layouts[5]  # Title Only
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = item["name"]

    left = Inches(0.5)
    top = Inches(1.2)
    width = Inches(5.5)

    tf_box = slide.shapes.add_textbox(left, top, width, Inches(4.5))
    tf = tf_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = item["description"]
    p.font.size = Pt(18)

    for point in item["anesthesia"]:
        bullet = tf.add_paragraph()
        bullet.text = f"- {point}"
        bullet.level = 1
        bullet.font.size = Pt(18)

    if image_path and os.path.exists(image_path):
        pic_left = Inches(6.2)
        pic_top = Inches(1.2)
        pic_width = Inches(3.0)
        slide.shapes.add_picture(image_path, pic_left, pic_top, width=pic_width)
        cap_tf = slide.shapes.add_textbox(pic_left, pic_top + Inches(3.05), pic_width, Inches(1.0)).text_frame
        cap_tf.text = item.get("image_caption", "")
        cap_tf.paragraphs[0].font.size = Pt(10)


def add_references_slide(prs: Presentation, refs: list[str]):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "References"
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for ref in refs:
        p = tf.add_paragraph()
        p.text = ref
        p.font.size = Pt(12)
        p.level = 0


def fetch_wikipedia_thumbnail(page_title: str) -> tuple[str | None, str | None]:
    """Return (thumb_url, page_url) for a Wikipedia page title, or (None, None)."""
    try:
        endpoint = "https://en.wikipedia.org/w/api.php"
        params = {
            "action": "query",
            "format": "json",
            "redirects": 1,
            "titles": page_title,
            "prop": "pageimages",
            "pithumbsize": 800,
        }
        headers = {
            "User-Agent": (
                "Mozilla/5.0 (X11; Linux x86_64) AppleWebKit/537.36 "
                "(KHTML, like Gecko) Chrome/123.0 Safari/537.36"
            ),
        }
        r = requests.get(endpoint, params=params, headers=headers, timeout=30)
        r.raise_for_status()
        data = r.json()
        pages = data.get("query", {}).get("pages", {})
        for page_id, page in pages.items():
            if "thumbnail" in page:
                thumb = page["thumbnail"]["source"]
                title = page.get("title", page_title)
                page_url = f"https://en.wikipedia.org/wiki/{urllib.parse.quote(title.replace(' ', '_'))}"
                return thumb, page_url
        return None, None
    except Exception as e:
        print(f"Wikipedia API failed for {page_title}: {e}")
        return None, None


def main():
    prs = Presentation()
    add_title_slide(
        prs,
        "Pediatric Craniofacial Syndromes",
        "Key types, features, and anesthesia considerations"
    )

    refs = []

    fallback_titles = {
        "Apert Syndrome": "Apert syndrome",
        "Crouzon Syndrome": "Crouzon syndrome",
        "Pfeiffer Syndrome": "Pfeiffer syndrome",
        "Saethre-Chotzen Syndrome": "Saethre–Chotzen syndrome",
        "Carpenter Syndrome": "Carpenter syndrome",
        "Muenke Syndrome": "Muenke syndrome",
        "Treacher Collins Syndrome": "Treacher Collins syndrome",
        "Pierre Robin Sequence": "Pierre Robin sequence",
        "Hemifacial Microsomia / Goldenhar": "Hemifacial microsomia",
        "Freeman-Sheldon Syndrome": "Freeman–Sheldon syndrome",
    }

    for item in syndromes:
        filename = item["name"].lower().replace(" ", "_").replace("/", "-") + ".jpg"
        dest_path = os.path.join(IMAGES_DIR, filename)

        img_path = None
        img_url_used = None

        if item.get("image_url"):
            img_path = download_image(item["image_url"], dest_path)
            img_url_used = item["image_url"] if img_path else None

        if not img_path:
            title = fallback_titles.get(item["name"], item["name"])            
            thumb_url, page_url = fetch_wikipedia_thumbnail(title)
            if thumb_url:
                img_path = download_image(thumb_url, dest_path)
                img_url_used = page_url or thumb_url

        if img_path and img_url_used:
            refs.append(f"Image: {item['name']} — {img_url_used}")

        add_syndrome_slide(prs, item, img_path)

    refs.extend([
        "General references: IntechOpen, Medscape, NCBI Bookshelf, ClinicalGate, Aneskey",
        "Always verify local protocols and latest guidelines"
    ])
    add_references_slide(prs, refs)

    prs.save(OUTPUT_PPTX)
    print(f"Saved: {OUTPUT_PPTX}")

if __name__ == "__main__":
    main()
